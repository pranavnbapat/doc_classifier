from __future__ import annotations

from pptx import Presentation

from docint.ingest.models import IngestedAsset


def ingest_pptx(file_path: str, filename: str) -> IngestedAsset:
    pres = Presentation(file_path)
    lines = []

    for idx, slide in enumerate(pres.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    line = line.strip()
                    if line:
                        slide_lines.append(line)
        if slide_lines:
            lines.append(f"Slide {idx}")
            lines.extend(slide_lines)

    text = "\n".join(lines).strip()

    return IngestedAsset(
        asset_type="pptx",
        filename=filename,
        source_path=file_path,
        text=text,
        lines=lines,
        units=len(pres.slides),
        unit_label="slides",
        source="pptx_text",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        visual_candidate=True,
        ocr_supported=False,
        meta={"slide_count": len(pres.slides)},
    )
